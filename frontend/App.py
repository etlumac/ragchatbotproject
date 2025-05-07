import time
import asyncio
import docx
import os
import json
import streamlit as st
from pathlib import Path
from pptx import Presentation
import base64
import streamlit.components.v1 as components
from services.file_manager import index_document, delete_document, safe_filename
from services.supabase_utils import download_file, get_public_url, list_uploaded_files
from services.auth import login_user, register_user

import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "backend")))
from app.routes.chat import chat_with_rag
from app.routes.chat import ChatRequest
from dotenv import load_dotenv
load_dotenv()
from services.supabase_db import (
    fetch_metadata,
    fetch_chat_history,
    save_chat_history,
    delete_metadata,
    delete_chat_history
)

st.set_page_config(layout="wide", page_title="RAGchat")

CHAT_HISTORY_PATH = Path("static/chat_history.json")
METADATA_PATH = Path("static/metadata.json")

SUPPORTED_EXTENSIONS = ("pdf", "docx", "txt", "pptx")


def load_json(filepath):
    if filepath.exists():
        with open(filepath, "r", encoding="utf-8") as file:
            return json.load(file)
    return {}

def save_json(filepath, data):
    with open(filepath, "w", encoding="utf-8") as file:
        json.dump(data, file, ensure_ascii=False, indent=4)

# === АВТОРИЗАЦИЯ ПОЛЬЗОВАТЕЛЯ ===
if "user_id" not in st.session_state:
    st.sidebar.title("Вход / Регистрация")

    mode = st.sidebar.radio("Выберите режим", ["Вход", "Регистрация"])
    email = st.sidebar.text_input("Email")
    password = st.sidebar.text_input("Пароль", type="password")

    if mode == "Регистрация":
        username = st.sidebar.text_input("Имя пользователя")
        if st.sidebar.button("Зарегистрироваться"):
            try:
                register_user(username, email, password)
                st.success("Пользователь зарегистрирован. Теперь войдите.")
            except Exception as e:
                st.error(f"Ошибка регистрации: {e}")
    else:
        if st.sidebar.button("Войти"):
            user_id, error = login_user(email, password)
            if user_id:
                st.session_state["user_id"] = user_id
                st.success("Успешный вход!")
                st.rerun()
            else:
                st.error(error)

    st.stop()  # Остановить выполнение если не вошёл



# Initialize session state
user_id = st.session_state["user_id"]

if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = fetch_chat_history(user_id)

if "metadata" not in st.session_state:
    st.session_state["metadata"] = fetch_metadata(user_id)

if "documents" not in st.session_state:
    st.session_state["documents"] = list_uploaded_files(user_id)

col1, col_mid, col2 = st.columns([1, 0.05, 1])

def display_pdf(file_path):
    with open(file_path, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode("utf-8")
    pdf_display = f"""
        <iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800px" type="application/pdf"></iframe>
    """
    st.markdown(pdf_display, unsafe_allow_html=True)


# ==== КОЛОНКА 1: ЗАГРУЗКА ФАЙЛОВ ====
with col1:
    if st.sidebar.button("🚪 Выход"):
        del st.session_state["user_id"]
        st.rerun()

    st.sidebar.title("Загруженные файлы")

    if st.session_state.get("reset_uploader"):
        st.session_state["reset_uploader"] = False

    if "file_uploaded" not in st.session_state:
        st.session_state["file_uploaded"] = False

    uploader_key = f"file_uploader_{time.time()}" if st.session_state.get("reset_uploader") else "file_uploader"

    uploaded_file = st.sidebar.file_uploader(
    "Загрузите документ",
    type=list(SUPPORTED_EXTENSIONS),
    key=uploader_key
)

    if uploaded_file and uploaded_file.name not in st.session_state["documents"]:
        try:
            #result = index_document(uploaded_file)
            result = index_document(uploaded_file, user_id)
            if "error" in result:
                st.sidebar.error(f"❌ Ошибка: {result['error']}")
            else:
                st.sidebar.success(f"✅ Файл {uploaded_file.name} успешно проиндексирован!")

                st.session_state["documents"].append(uploaded_file.name)
                st.session_state["last_selected_document"] = uploaded_file.name
                st.session_state["file_uploaded"] = True

                # 👇 Сброс file_uploader после загрузки
                #st.experimental_set_query_params(file_uploaded="1")
                st.session_state["reset_uploader"] = True
                st.query_params["file_uploaded"] = "1"
                st.rerun()
        except Exception as e:
            st.sidebar.error(f"❌ Ошибка при загрузке: {e}")
    else:
        # 👇 Если флаг установлен в query params, сбрасываем его и file_uploader
        #if st.experimental_get_query_params().get("file_uploaded"):
        if 'file_uploaded' in st.query_params:
            st.session_state["file_uploaded"] = False
            #st.experimental_set_query_params()
            st.query_params.clear()

    selected_document = (
        st.sidebar.selectbox("Текущий документ:", st.session_state["documents"])
        if st.session_state["documents"]
        else None
    )

    if selected_document is not None:
        if "last_selected_document" not in st.session_state or selected_document != st.session_state["last_selected_document"]:
            st.session_state["messages"] = st.session_state["chat_history"].get(
                selected_document,
                [{"role": "ai", "content": "Привет! Я готов помочь вам с вашим документом."}]
            )
            st.session_state["last_selected_document"] = selected_document
    else:
        st.session_state["messages"] = [{"role": "ai", "content": "Привет! Загрузите новый документ, чтобы начать."}]

    # === Отображение содержимого выбранного документа ===
    if selected_document:
        #file_path = download_file(selected_document)
        file_path = download_file(user_id, selected_document)


        if file_path and file_path.exists():
            if selected_document.endswith(".pdf"):
                public_url = get_public_url(user_id, selected_document)
                viewer_url = f"https://docs.google.com/gview?url={public_url}&embedded=true"
                components.html(f'<iframe src="{viewer_url}" width="100%" height="800px" style="border: none;"></iframe>', height=800)

            elif selected_document.endswith(".docx"):
                doc = docx.Document(file_path)
                doc_text = "\n\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
                st.markdown(f'<div style="height:800px; overflow:auto; border:1px solid #ddd; padding:10px;">{doc_text}</div>', unsafe_allow_html=True)

            elif selected_document.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                st.markdown(f'<div style="height:800px; overflow:auto; border:1px solid #ddd; padding:10px;">{content}</div>', unsafe_allow_html=True)

            elif selected_document.endswith(".pptx"):
                prs = Presentation(file_path)
                ppt_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            ppt_text.append(shape.text.strip())
                full_text = "\n\n".join(ppt_text)
                st.markdown(f'<div style="height:800px; overflow:auto; border:1px solid #ddd; padding:10px;">{full_text}</div>', unsafe_allow_html=True)

# Функция удаления чата
def delete_chat():
    try:
        #delete_document(selected_document)
        delete_document(selected_document, user_id)

        safe_name = safe_filename(selected_document)

        st.session_state["chat_history"].pop(selected_document, None)
        #delete_chat_history(selected_document)
        delete_chat_history(user_id, safe_name)

        st.session_state["metadata"].pop(safe_name, None)
        #delete_metadata(safe_name)
        delete_metadata(user_id, safe_name)

        st.session_state["documents"].remove(selected_document)
        st.session_state["last_selected_document"] = st.session_state["documents"][0] if st.session_state["documents"] else None

        if not st.session_state["documents"]:
            st.session_state["messages"] = [{"role": "ai", "content": "Привет! Загрузите новый документ, чтобы начать."}]

         #st.success("✅ Чат и документ удалены!")
    except Exception as e:
        print(f"❌ Ошибка при удалении: {e}")
        st.error(f"❌ Ошибка при удалении: {e}")
        time.sleep(10)
    st.rerun()


# ==== КОЛОНКА 2: ЧАТ ====
with col2:
    if selected_document:
        col_buttons = st.columns([0.25, 0.57]) 

        with col_buttons[0]:
            if st.button("🔄 Сбросить историю чата"):
                if selected_document in st.session_state["chat_history"]:
                    del st.session_state["chat_history"][selected_document]
                    safe_name = safe_filename(selected_document)
                    save_chat_history(user_id, safe_name, [{"role": "ai", "content": "Привет! Я готов помочь вам с вашим документом."}])
                st.session_state["messages"] = [{"role": "ai", "content": "Привет! Я готов помочь вам с вашим документом."}]
                st.rerun()

        with col_buttons[1]:
            if st.button("🗑️ Удалить чат и документ", key="delete_chat_button"):
                delete_chat()


        chat_container = st.container()
        with chat_container:
            chat_box = st.container(height=650)
            with chat_box:
                for msg in st.session_state["messages"]:
                    with st.chat_message(msg["role"]):
                        st.write(msg["content"])

        user_input = st.chat_input("Введите ваш вопрос")

        if user_input:
            st.session_state["messages"].append({"role": "user", "content": user_input})
            with chat_box:
                with st.chat_message("user"):
                    st.write(user_input)

            with st.spinner("Обрабатываю..."):
                try:
                    #chat_request = ChatRequest(query=user_input, document_name=selected_document)
                    #response = asyncio.run(chat_with_rag(chat_request))

                    chat_request = ChatRequest(query=user_input, document_name=selected_document, user_id=user_id)
                    response = asyncio.run(chat_with_rag(chat_request))

                    assistant_response = response.get("response", "Ошибка: ответ не найден.")
                except Exception as e:
                    assistant_response = f"❌ Ошибка обработки запроса: {e}"

            st.session_state["messages"].append({"role": "ai", "content": assistant_response})
            with chat_box:
                with st.chat_message("ai"):
                    st.write(assistant_response)

            st.session_state["chat_history"][selected_document] = st.session_state["messages"]
            #save_chat_history(selected_document, st.session_state["messages"])
            safe_name = safe_filename(selected_document)
            save_chat_history(user_id, safe_name, st.session_state["messages"])

    #else:
        #st.info("Загрузите документ, чтобы начать чат.")


