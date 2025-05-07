import json
from pathlib import Path
import docx
from pptx import Presentation
from langchain.docstore.document import Document as LangchainDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from dotenv import load_dotenv
from slugify import slugify
from services.supabase_utils import upload_file, delete_file
from services.qdrant_utils import upload_vectors, delete_vectors_by_document_id
from services.qdrant_utils import init_qdrant
from services.supabase_db import save_metadata as save_metadata_to_db
import os
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")))
from nltk.tokenize import sent_tokenize
import nltk
nltk.download("punkt")
nltk.download('punkt_tab')
try:
    nltk.data.find("tokenizers/punkt")
except LookupError:
    nltk.download("punkt")
    nltk.download('punkt_tab')

load_dotenv()

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".txt", ".pptx")
METADATA_PATH = Path("static/metadata.json")

def sentence_window_split(text, window_size=3, stride=1):
    sentences = sent_tokenize(text)
    chunks = []
    for i in range(0, len(sentences) - window_size + 1, stride):
        chunk_text = " ".join(sentences[i:i+window_size])
        chunks.append(LangchainDocument(page_content=chunk_text))
    return chunks

def safe_filename(filename: str) -> str:
    name, ext = Path(filename).stem, Path(filename).suffix
    return f"{slugify(name)}{ext}"

def load_metadata():
    if METADATA_PATH.exists():
        with open(METADATA_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def save_metadata(metadata):
    with open(METADATA_PATH, "w", encoding="utf-8") as f:
        json.dump(metadata, f, indent=4, ensure_ascii=False)

def extract_text_from_file(file_path: Path) -> str:
    if file_path.suffix == ".pdf":
        import fitz
        return "\n".join([page.get_text() for page in fitz.open(file_path)])
    elif file_path.suffix == ".docx":
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_path.suffix == ".txt":
        return file_path.read_text(encoding="utf-8")
    elif file_path.suffix == ".pptx":
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    else:
        raise ValueError("Unsupported file format")

def index_document(uploaded_file, user_id: str):
    filename = uploaded_file.name
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return {"error": "Формат не поддерживается"}

    temp_path = Path(f"temp_{filename}")
    with open(temp_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    try:
        text = extract_text_from_file(temp_path)

        raw_doc = LangchainDocument(page_content=text, metadata={"source": filename})
        chunks = sentence_window_split(raw_doc.page_content, window_size=5, stride=1)
        if not chunks:
            return {"error": "Не удалось разбить документ на чанки"}

        embeddings = OpenAIEmbeddings(
            model="emb-openai/text-embedding-3-small",
            openai_api_key=os.getenv("OPENAI_API_KEY"),
            openai_api_base="https://api.vsegpt.ru/v1"
        )

        texts = [chunk.page_content for chunk in chunks]
        metadatas = [
            {
                **chunk.metadata,
                "text": chunk.page_content
            }
            for chunk in chunks
        ]
        vectors = embeddings.embed_documents(texts)

        doc_id = safe_filename(filename)

        upload_vectors(user_id, doc_id, vectors, metadatas)
        save_metadata_to_db(user_id, doc_id, {"status": "uploaded"})
        upload_file(user_id, filename, temp_path)

        temp_path.unlink(missing_ok=True)

        return {"message": "Файл проиндексирован", "filename": filename}
    
    except Exception as e:
        return {"error": f"Ошибка при индексации: {e}"}

def delete_document(filename: str, user_id: str):
    delete_file(user_id, filename)

    doc_id = safe_filename(filename)
    delete_vectors_by_document_id(doc_id, user_id)


    metadata = load_metadata()
    if doc_id in metadata:
        del metadata[doc_id]
        save_metadata(metadata)

